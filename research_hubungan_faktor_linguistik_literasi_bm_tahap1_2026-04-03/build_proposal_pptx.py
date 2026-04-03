from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
OUT_DIR = BASE / "pptx"
OUT_DIR.mkdir(exist_ok=True)
OUTPUT = OUT_DIR / "06_pembentangan_proposal.pptx"

TITLE = (
    "Hubungan antara Faktor Linguistik dan Keberkesanan Strategi "
    "Pengajaran Literasi Bahasa Melayu Murid Tahap 1"
)


COLORS = {
    "navy": RGBColor(14, 77, 100),
    "teal": RGBColor(31, 122, 140),
    "gold": RGBColor(217, 164, 65),
    "sand": RGBColor(245, 240, 231),
    "cream": RGBColor(252, 250, 246),
    "ink": RGBColor(33, 43, 54),
    "muted": RGBColor(93, 108, 122),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(215, 221, 227),
    "green": RGBColor(86, 148, 95),
    "red": RGBColor(175, 74, 74),
}


def set_bg(slide, color_key="cream"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[color_key]


def add_top_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["navy"]
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(10.6), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS["white"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(0.55), Inches(9.2), Inches(0.2))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(214, 228, 232)


def add_footer(slide, text="Cadangan pembentangan proposal kajian"):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.32)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["navy"]
    line.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.45), Inches(7.2), Inches(9.5), Inches(0.18))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["white"]


def add_bullets(slide, left, top, width, height, items, font_size=20, color_key="ink"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[color_key]
        p.space_after = Pt(8)


def add_body_text(slide, left, top, width, height, text, font_size=18, color_key="ink", bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color_key]
    return box


def add_section_label(slide, text, left, top, width=Inches(2.4)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.38)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["gold"]
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]


def add_card(slide, left, top, width, height, title, body, accent="teal"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1.1)

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS[accent]
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.12), width - Inches(0.3), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS["navy"]

    body_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.46), width - Inches(0.32), height - Inches(0.52))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(12.5)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(4)


def add_stat_card(slide, left, top, width, height, number, label, note, accent="teal"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)

    num = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.42))
    p = num.text_frame.paragraphs[0]
    p.text = number
    p.font.name = "Aptos Display"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS[accent]

    lab = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.55), width - Inches(0.3), Inches(0.34))
    p = lab.text_frame.paragraphs[0]
    p.text = label
    p.font.name = "Aptos"
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]

    note_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.95), width - Inches(0.3), Inches(0.46))
    p = note_box.text_frame.paragraphs[0]
    p.text = note
    p.font.name = "Aptos"
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLORS["muted"]


def add_reference_line(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.72), Inches(12.0), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(8.5)
    p.font.italic = True
    p.font.color.rgb = COLORS["muted"]


def add_chart(slide):
    chart_data = CategoryChartData()
    chart_data.categories = ["Disaring", "Belum kuasai 3M", "Berjaya selepas intervensi"]
    chart_data.add_series("Bilangan murid", (448113, 122062, 48308))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(1.65), Inches(5.5), Inches(3.5), chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = 500000
    chart.value_axis.minimum_scale = 0
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["teal"]


def add_framework_diagram(slide):
    boxes = [
        (Inches(0.9), Inches(2.05), "Kesedaran\nfonologi"),
        (Inches(0.9), Inches(3.05), "Kesedaran\nmorfologi"),
        (Inches(0.9), Inches(4.05), "Kosa kata &\nbahasa lisan"),
        (Inches(0.9), Inches(5.05), "Bahasa ibunda /\nbahasa rumah"),
    ]

    target = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(3.15), Inches(2.8), Inches(1.2)
    )
    target.fill.solid()
    target.fill.fore_color.rgb = COLORS["teal"]
    target.line.fill.background()
    p = target.text_frame.paragraphs[0]
    p.text = "Keberkesanan\nstrategi pengajaran"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["white"]

    outcome = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.35), Inches(3.15), Inches(2.8), Inches(1.2)
    )
    outcome.fill.solid()
    outcome.fill.fore_color.rgb = COLORS["gold"]
    outcome.line.fill.background()
    p = outcome.text_frame.paragraphs[0]
    p.text = "Hasil literasi murid\nbacaan, kefahaman, tulisan"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["white"]

    for left, top, label in boxes:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(2.65), Inches(0.72)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["white"]
        shape.line.color.rgb = COLORS["line"]
        p = shape.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["ink"]

        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, left + Inches(2.65), top + Inches(0.36), Inches(5.35), Inches(3.75)
        )
        connector.line.color.rgb = COLORS["teal"]
        connector.line.width = Pt(1.5)

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(8.15), Inches(3.75), Inches(9.35), Inches(3.75)
    )
    connector.line.color.rgb = COLORS["gold"]
    connector.line.width = Pt(2)


def add_timeline_table(slide):
    table = slide.shapes.add_table(10, 3, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.7)).table
    table.columns[0].width = Inches(1.1)
    table.columns[1].width = Inches(7.2)
    table.columns[2].width = Inches(3.6)

    headers = ["Fasa", "Aktiviti", "Tempoh"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Aptos"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

    rows = [
        ("1", "Pembinaan proposal, semakan literatur, pembangunan instrumen", "Bulan 1-2"),
        ("2", "Semakan pakar, kajian rintis, penambahbaikan instrumen", "Bulan 3"),
        ("3", "Kebenaran dan penyelarasan dengan sekolah", "Bulan 4"),
        ("4", "Ujian pra dan pengumpulan data kuantitatif", "Bulan 5"),
        ("5", "Pemerhatian PdP dan tempoh pengajaran", "Bulan 5-6"),
        ("6", "Ujian pasca dan analisis awal", "Bulan 7"),
        ("7", "Temu bual dan pemerhatian kualitatif", "Bulan 8"),
        ("8", "Integrasi dapatan dan interpretasi", "Bulan 9"),
        ("9", "Penulisan laporan akhir", "Bulan 10"),
    ]

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"] if r % 2 else COLORS["sand"]
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["ink"]
            if c != 1:
                p.alignment = PP_ALIGN.CENTER


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide, "sand")
    left_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(4.55), Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLORS["navy"]
    left_panel.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.35), Inches(0), Inches(0.2), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["gold"]
    accent.line.fill.background()
    add_section_label(slide, "PEMBENTANGAN PROPOSAL", Inches(5.1), Inches(0.7), Inches(3.0))
    add_body_text(slide, Inches(5.1), Inches(1.45), Inches(7.5), Inches(1.8), TITLE, 24, "ink", True)
    add_body_text(
        slide,
        Inches(5.1),
        Inches(3.25),
        Inches(6.6),
        Inches(1.0),
        "Cadangan kajian untuk meneliti padanan antara profil linguistik murid Tahap 1 dengan strategi pengajaran literasi Bahasa Melayu.",
        18,
        "muted",
    )
    add_body_text(slide, Inches(0.6), Inches(1.3), Inches(3.2), Inches(0.4), "Nama Penyelidik", 14, "white", True)
    add_body_text(slide, Inches(0.6), Inches(1.75), Inches(3.2), Inches(0.4), "______________________________", 14, "white")
    add_body_text(slide, Inches(0.6), Inches(2.45), Inches(3.2), Inches(0.4), "Program / Fakulti", 14, "white", True)
    add_body_text(slide, Inches(0.6), Inches(2.9), Inches(3.2), Inches(0.4), "______________________________", 14, "white")
    add_body_text(slide, Inches(0.6), Inches(5.95), Inches(3.1), Inches(0.5), "3 April 2026", 16, "white", True)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Struktur Pembentangan", "Aliran ringkas proposal kajian")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(1.45),
        Inches(5.1),
        Inches(4.8),
        [
            "Latar belakang, isu dan keperluan kajian",
            "Objektif, soalan kajian dan hipotesis",
            "Sorotan literatur utama dan jurang penyelidikan",
            "Kerangka konseptual hubungan faktor linguistik dengan strategi pengajaran",
            "Metodologi, sampel, instrumen dan analisis data",
            "Jadual pelaksanaan dan sumbangan kajian",
        ],
        19,
    )
    add_card(
        slide,
        Inches(7.0),
        Inches(1.65),
        Inches(5.2),
        Inches(3.35),
        "Fokus teras pembentangan",
        "Kajian ini menguji sama ada strategi pengajaran literasi menjadi lebih berkesan apabila dipadankan dengan keperluan linguistik murid.\n\nCadangan ini menumpukan murid Tahap 1 kerana fasa ini kritikal untuk pembinaan asas membaca, menulis dan kefahaman.",
        "gold",
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Latar Belakang Isu", "Data rasmi menunjukkan keperluan intervensi literasi awal yang lebih tepat")
    add_stat_card(slide, Inches(0.8), Inches(1.45), Inches(2.2), Inches(1.55), "448,113", "Murid Tahun 1 disaring", "Pengesanan awal literasi dan numerasi", "teal")
    add_stat_card(slide, Inches(3.15), Inches(1.45), Inches(2.2), Inches(1.55), "122,062", "Belum kuasai 3M", "Sekitar 27.2% daripada jumlah disaring", "red")
    add_stat_card(slide, Inches(5.5), Inches(1.45), Inches(2.2), Inches(1.55), "48,308", "Berjaya selepas intervensi", "Sekitar 39.6% daripada kumpulan sasaran", "green")
    add_chart(slide)
    add_body_text(
        slide,
        Inches(0.85),
        Inches(3.45),
        Inches(5.6),
        Inches(1.85),
        "Isyarat utama daripada data ini ialah masalah literasi awal masih besar, tetapi intervensi yang tepat mampu menghasilkan peningkatan yang bermakna dalam tempoh singkat. Soalan pentingnya ialah: strategi mana paling sesuai untuk profil linguistik murid yang berbeza?",
        16,
    )
    add_reference_line(
        slide,
        "Sumber: Kementerian Pendidikan Malaysia, Perhimpunan Bulanan KPM, 19 Disember 2024; Manual Pengoperasian Pengesanan Literasi dan Numerasi serta Program Intervensi Tahun 1 (2024).",
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Pernyataan Masalah", "Jurang antara keperluan murid dengan respons pedagogi")
    add_bullets(
        slide,
        Inches(0.85),
        Inches(1.45),
        Inches(6.1),
        Inches(4.6),
        [
            "Sebahagian murid Tahap 1 masih belum menguasai literasi asas walaupun intervensi telah dilaksanakan.",
            "Masalah murid tidak seragam: ada yang lemah pada fonologi, ada yang terbatas dari segi kosa kata, morfologi atau bahasa rumah.",
            "Banyak kajian tempatan membincangkan masalah murid atau strategi guru secara berasingan, tetapi kurang menguji hubungan langsung antara kedua-duanya.",
            "Tanpa data tentang padanan strategi dengan faktor linguistik, intervensi berisiko kekal umum dan kurang tepat sasaran.",
        ],
        18,
    )
    add_card(
        slide,
        Inches(7.35),
        Inches(1.65),
        Inches(4.95),
        Inches(3.8),
        "Soalan utama kajian",
        "Adakah strategi pengajaran literasi Bahasa Melayu menjadi lebih berkesan apabila disesuaikan dengan profil linguistik murid Tahap 1?\n\nJika ya, faktor linguistik manakah yang paling kuat mempengaruhi keberkesanannya?",
        "teal",
    )
    add_reference_line(
        slide,
        "Rujukan utama: Chew (2016); Wan Ahmad (2019); Nahar (2020); Low & Mahamod (2024); Low et al. (2024)."
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Tujuan, Objektif dan Soalan Kajian")
    add_section_label(slide, "Tujuan", Inches(0.8), Inches(1.2), Inches(1.3))
    add_body_text(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(5.6),
        Inches(1.15),
        "Meneliti hubungan antara faktor linguistik murid Tahap 1 dengan keberkesanan strategi pengajaran literasi Bahasa Melayu di sekolah rendah.",
        18,
    )
    add_section_label(slide, "Objektif", Inches(0.8), Inches(2.7), Inches(1.5))
    add_bullets(
        slide,
        Inches(0.8),
        Inches(3.1),
        Inches(5.8),
        Inches(3.0),
        [
            "Mengenal pasti tahap faktor linguistik utama murid.",
            "Mengenal pasti strategi pengajaran yang digunakan guru.",
            "Menentukan hubungan antara faktor linguistik dengan keberkesanan strategi.",
            "Mengenal pasti faktor linguistik yang paling dominan meramal keberkesanan.",
            "Meneroka penyesuaian strategi berdasarkan profil murid.",
        ],
        15,
    )
    add_card(
        slide,
        Inches(7.05),
        Inches(1.45),
        Inches(5.15),
        Inches(4.85),
        "Soalan kajian teras",
        "1. Apakah tahap kesedaran fonologi, morfologi, kosa kata dan pengaruh bahasa rumah murid?\n2. Apakah strategi pengajaran yang paling kerap digunakan guru?\n3. Adakah terdapat hubungan signifikan antara faktor linguistik dengan keberkesanan strategi?\n4. Faktor manakah paling dominan meramal keberkesanan?\n5. Bagaimanakah guru menyesuaikan strategi mengikut profil murid?",
        "gold",
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Sorotan Literatur: Faktor Linguistik")
    add_card(slide, Inches(0.8), Inches(1.35), Inches(2.85), Inches(2.0), "Kesedaran fonologi", "Asas kepada hubungan bunyi-huruf dan penyahkodan.\nMurid yang lemah pada bunyi, rima dan suku kata cenderung lebih sukar membaca awal.", "teal")
    add_card(slide, Inches(3.85), Inches(1.35), Inches(2.85), Inches(2.0), "Kesedaran morfologi", "Penting untuk memahami kata dasar dan kata terbitan.\nMakin penting apabila teks semakin kaya dengan imbuhan.", "gold")
    add_card(slide, Inches(6.9), Inches(1.35), Inches(2.85), Inches(2.0), "Kosa kata & bahasa lisan", "Menentukan kefahaman, bukan sekadar kebolehan menyebut.\nPendedahan bahasa yang kaya membantu pembacaan bermakna.", "green")
    add_card(slide, Inches(9.95), Inches(1.35), Inches(2.55), Inches(2.0), "Bahasa rumah", "Bahasa ibunda dan ekologi bahasa memberi kesan pada input, sebutan dan keyakinan menggunakan BM.", "red")
    add_body_text(
        slide,
        Inches(0.9),
        Inches(4.0),
        Inches(11.5),
        Inches(1.2),
        "Sorotan literatur menunjukkan prestasi literasi awal tidak boleh dijelaskan melalui satu faktor sahaja. Sebaliknya, prestasi murid terbentuk melalui gabungan komponen fonologi, morfologi, kosa kata dan pengalaman bahasa harian.",
        18,
    )
    add_reference_line(
        slide,
        "Rujukan utama: Hoover & Gough (1990); Sumardi (2010); Chin et al. (2012); Nahar (2020); Low et al. (2024)."
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Sorotan Literatur: Strategi Pengajaran dan Jurang Kajian")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.45),
        Inches(5.8),
        Inches(4.6),
        [
            "Pengajaran langsung dan eksplisit penting untuk bunyi-huruf, suku kata dan bacaan awal.",
            "Pengajaran kumpulan kecil dan bimbingan berfokus membantu murid yang belum menguasai konstruk tertentu.",
            "Analisis kosa kata, bahan bacaan terkawal dan aktiviti berpusatkan murid meningkatkan penglibatan dan kefahaman.",
            "Dalam konteks SJKC, guru sering menyesuaikan strategi dengan cabaran bahasa ibunda dan minat murid.",
        ],
        18,
    )
    add_card(
        slide,
        Inches(7.2),
        Inches(1.55),
        Inches(5.0),
        Inches(3.9),
        "Jurang penyelidikan",
        "Kajian tempatan banyak meneroka masalah penguasaan murid atau amalan guru secara berasingan.\n\nMasih terhad kajian yang menguji hubungan langsung antara faktor linguistik murid Tahap 1 dengan keberkesanan strategi pengajaran dalam satu model kajian yang menyeluruh.",
        "teal",
    )
    add_reference_line(
        slide,
        "Rujukan utama: IES (2017); Wan Ahmad (2019); Hamzah & Mahamod (2021); Low & Mahamod (2024)."
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Kerangka Konseptual Kajian", "Strategi berkesan apabila dipadankan dengan profil linguistik murid")
    add_framework_diagram(slide)
    add_body_text(
        slide,
        Inches(0.95),
        Inches(6.05),
        Inches(11.4),
        Inches(0.46),
        "Andaian utama: semakin tepat guru memadankan strategi pengajaran dengan keperluan linguistik murid, semakin tinggi kebarangkalian peningkatan literasi berlaku.",
        14.5,
        "muted",
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Metodologi Kajian", "Reka bentuk kaedah campuran explanatory sequential")
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.6), Inches(3.55), "Fasa 1: Kuantitatif", "Ujian faktor linguistik\nSoal selidik strategi guru\nUjian literasi pra dan pasca\nAnalisis deskriptif, korelasi dan regresi", "teal")
    add_card(slide, Inches(4.85), Inches(1.5), Inches(3.6), Inches(3.55), "Fasa 2: Kualitatif", "Temu bual guru terpilih\nPemerhatian bilik darjah\nAnalisis tematik\nPenjelasan terhadap pola dapatan kuantitatif", "gold")
    add_card(slide, Inches(8.9), Inches(1.5), Inches(3.6), Inches(3.55), "Integrasi", "Meta-inferens tentang hubungan faktor linguistik dengan keberkesanan strategi pengajaran\nCadangan amalan intervensi yang lebih bersasar", "green")
    arrow1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.4), Inches(3.25), Inches(4.85), Inches(3.25))
    arrow1.line.color.rgb = COLORS["navy"]
    arrow1.line.width = Pt(2)
    arrow2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.45), Inches(3.25), Inches(8.9), Inches(3.25))
    arrow2.line.color.rgb = COLORS["navy"]
    arrow2.line.width = Pt(2)
    add_body_text(
        slide,
        Inches(0.9),
        Inches(5.45),
        Inches(11.4),
        Inches(0.7),
        "Reka bentuk ini dipilih kerana sesuai untuk mengukur kekuatan hubungan terlebih dahulu, kemudian menjelaskan mengapa dan bagaimana strategi tertentu lebih berkesan untuk kumpulan murid yang berbeza.",
        15.5,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Populasi, Sampel dan Instrumen")
    add_card(slide, Inches(0.8), Inches(1.45), Inches(3.3), Inches(2.15), "Populasi & lokasi", "Sekolah rendah kerajaan dan bantuan kerajaan dalam sebuah daerah yang mempunyai kepelbagaian latar bahasa.\nFokus pada murid Tahap 1 dan guru Bahasa Melayu Tahap 1.", "teal")
    add_card(slide, Inches(4.35), Inches(1.45), Inches(3.3), Inches(2.15), "Sampel kuantitatif", "300 orang murid Tahap 1\n30 orang guru Bahasa Melayu\n10 buah sekolah", "gold")
    add_card(slide, Inches(7.9), Inches(1.45), Inches(4.35), Inches(2.15), "Sampel kualitatif", "6 orang guru dipilih secara bertujuan berdasarkan variasi dapatan.\n6 sesi pemerhatian bilik darjah untuk triangulasi.", "green")
    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.1),
        Inches(11.2),
        Inches(2.2),
        [
            "Ujian kesedaran fonologi",
            "Ujian kesedaran morfologi asas",
            "Ujian kosa kata dan bahasa lisan",
            "Soal selidik bahasa rumah",
            "Soal selidik strategi pengajaran guru",
            "Ujian literasi Bahasa Melayu pra dan pasca",
            "Borang pemerhatian PdP dan protokol temu bual guru",
        ],
        16,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Analisis Data dan Hipotesis")
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.55), Inches(3.55), "Hipotesis", "H1 hingga H4 menguji hubungan antara faktor linguistik dengan keberkesanan strategi.\nH5 menguji faktor peramal dominan terhadap keberkesanan strategi pengajaran.", "teal")
    add_card(slide, Inches(4.7), Inches(1.5), Inches(3.55), Inches(3.55), "Analisis kuantitatif", "Statistik deskriptif\nUjian normaliti\nKorelasi Pearson\nRegresi berganda\nPaired sample t-test\nANOVA jika perlu", "gold")
    add_card(slide, Inches(8.6), Inches(1.5), Inches(3.7), Inches(3.55), "Analisis kualitatif", "Pengekodan terbuka\nPengelompokan kategori\nPembinaan tema\nIntegrasi dapatan dengan hasil kuantitatif", "green")
    add_body_text(
        slide,
        Inches(0.9),
        Inches(5.45),
        Inches(11.4),
        Inches(0.68),
        "Fokus analitik bukan sekadar melihat sama ada strategi digunakan, tetapi sejauh mana strategi itu benar-benar berkesan untuk murid dengan profil linguistik yang berbeza.",
        15.5,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Jadual Pelaksanaan Kajian")
    add_timeline_table(slide)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Sumbangan Kajian")
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(3.7), "Sumbangan teori", "Memperjelas hubungan antara faktor linguistik murid dengan keberkesanan strategi pengajaran literasi Bahasa Melayu dalam konteks Tahap 1.", "teal")
    add_card(slide, Inches(4.85), Inches(1.5), Inches(3.7), Inches(3.7), "Sumbangan amalan", "Membantu guru memilih strategi yang lebih tepat berdasarkan masalah fonologi, morfologi, kosa kata atau bahasa rumah murid.", "gold")
    add_card(slide, Inches(8.9), Inches(1.5), Inches(3.4), Inches(3.7), "Sumbangan dasar", "Menyediakan evidens tempatan untuk menyokong intervensi awal, saringan dan panduan pedagogi yang lebih sensitif terhadap kepelbagaian bahasa murid.", "green")
    add_body_text(
        slide,
        Inches(0.9),
        Inches(5.6),
        Inches(11.3),
        Inches(0.6),
        "Hasil akhir yang diharapkan ialah model padanan strategi pengajaran dengan profil linguistik murid Tahap 1 yang boleh digunakan oleh guru dan pihak sekolah.",
        15.5,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide, "sand")
    add_top_bar(slide, "Rujukan Terpilih", "APA 7 ringkas untuk pembentangan")
    add_bullets(
        slide,
        Inches(0.85),
        Inches(1.4),
        Inches(12.0),
        Inches(4.9),
        [
            "Chew, F. P. (2016). Masalah pembelajaran Bahasa Melayu dalam kalangan murid Cina sekolah rendah. Jurnal Pendidikan Bahasa Melayu, 6(2), 10-22.",
            "Hoover, W. A., & Gough, P. B. (1990). The simple view of reading. Reading and Writing: An Interdisciplinary Journal, 2(2), 127-160.",
            "Kementerian Pendidikan Malaysia. (2016). Dokumen Standard Kurikulum dan Pentaksiran Bahasa Melayu Tahun 1 Sekolah Kebangsaan.",
            "Kementerian Pendidikan Malaysia. (2024a). Manual pengoperasian pengesanan literasi dan numerasi serta program intervensi Tahun 1.",
            "Low, J. Y., & Mahamod, Z. (2024). Strategi dan cabaran pengajaran kemahiran membaca... Jurnal Pendidikan Bahasa Melayu, 14(2), 117-137.",
            "Low, J. Y., Wee, X. H., Khoo, Y. T., & Anuaruddin, N. F. F. (2024). Persepsi guru terhadap pengaruh bahasa ibunda... Jurnal Pendidikan Bahasa Melayu, 14(1), 1-18.",
            "Nahar, N. (2020). Penguasaan kemahiran membaca dan menulis Bahasa Melayu... Issues in Language Studies, 9(1), 107-123.",
            "Wan Ahmad, W. N. (2019). Amalan guru dalam melaksanakan kemahiran literasi Bahasa Melayu bagi Program LINUS di sekolah rendah. Jurnal Pendidikan Bahasa Melayu, 9(1), 1-11.",
        ],
        12,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_top_bar(slide, "Penutup", "Terima kasih")
    add_body_text(
        slide,
        Inches(1.1),
        Inches(1.8),
        Inches(11.0),
        Inches(1.2),
        "Cadangan kajian ini berhasrat membantu kita bergerak daripada intervensi umum kepada intervensi literasi yang lebih tepat, responsif dan berasaskan profil linguistik murid.",
        24,
        "ink",
        True,
    )
    add_body_text(
        slide,
        Inches(1.1),
        Inches(3.35),
        Inches(8.5),
        Inches(1.0),
        "Sesi soal jawab dan maklum balas penyelia",
        20,
        "teal",
    )
    closing = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.3), Inches(2.4), Inches(2.4)
    )
    closing.fill.solid()
    closing.fill.fore_color.rgb = COLORS["gold"]
    closing.line.fill.background()
    p = closing.text_frame.paragraphs[0]
    p.text = "Q&A"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    add_footer(slide)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = build_presentation()
    print(output)
